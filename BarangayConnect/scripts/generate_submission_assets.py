from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape
from reportlab.lib.pagesizes import TABLOID
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
JPG_PATH = DOCS / "Final-ERD.jpg"
PDF_PATH = DOCS / "Final-ERD.pdf"
PPTX_PATH = DOCS / "IT114L-Project-Presentation.pptx"


BOXES = [
    {
        "title": "Residents",
        "lines": [
            "PK ResidentId",
            "FullName",
            "HouseholdNo",
            "ContactNumber",
            "EmailAddress",
            "Purok",
            "RegisteredOn",
        ],
        "xy": (80, 150),
    },
    {
        "title": "Services",
        "lines": [
            "PK ServiceId",
            "Name",
            "Office",
            "Description",
            "Schedule",
            "Requirements",
        ],
        "xy": (520, 150),
    },
    {
        "title": "Announcements",
        "lines": [
            "PK AnnouncementId",
            "Title",
            "Category",
            "Summary",
            "PublishedOn",
            "Audience",
        ],
        "xy": (980, 150),
    },
    {
        "title": "Appointments",
        "lines": [
            "PK AppointmentId",
            "FK ResidentId",
            "FK ServiceId",
            "AppointmentDate",
            "TimeSlot",
            "Status",
            "Notes",
        ],
        "xy": (210, 510),
    },
    {
        "title": "ServiceRequests",
        "lines": [
            "PK RequestId",
            "FK ResidentId",
            "FK ServiceId",
            "Description",
            "Priority",
            "Status",
            "SubmittedOn",
        ],
        "xy": (700, 510),
    },
]


def get_font(size: int, bold: bool = False):
    candidates = [
        "C:/Windows/Fonts/seguibl.ttf" if bold else "C:/Windows/Fonts/segoeui.ttf",
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
    ]
    for item in candidates:
        font_path = Path(item)
        if font_path.exists():
            return ImageFont.truetype(str(font_path), size)
    return ImageFont.load_default()


def draw_box(draw: ImageDraw.ImageDraw, title_font, body_font, title, lines, x, y):
    width = 320
    line_height = 28
    height = 80 + (len(lines) * line_height)
    draw.rounded_rectangle((x, y, x + width, y + height), radius=18, fill="white", outline="#1c2b2d", width=3)
    draw.rounded_rectangle((x, y, x + width, y + 54), radius=18, fill="#0c7c59", outline="#0c7c59")
    draw.text((x + 18, y + 12), title, fill="white", font=title_font)
    current_y = y + 68
    for line in lines:
        draw.text((x + 18, current_y), line, fill="#1c2b2d", font=body_font)
        current_y += line_height
    return width, height


def arrow(draw, start, end, color="#7a5412"):
    draw.line((start, end), fill=color, width=5)
    ex, ey = end
    draw.polygon([(ex, ey), (ex - 18, ey - 8), (ex - 18, ey + 8)], fill=color)


def generate_erd_jpg():
    image = Image.new("RGB", (1450, 920), "#f6f1e8")
    draw = ImageDraw.Draw(image)
    title_font = get_font(32, bold=True)
    section_font = get_font(18)
    body_font = get_font(20)

    draw.text((60, 40), "BarangayConnect Final Entity Relationship Diagram", fill="#143128", font=title_font)
    draw.text((62, 92), "Local government service portal using an embedded SQLite database", fill="#5e6e68", font=section_font)

    dimensions = {}
    for box in BOXES:
        dimensions[box["title"]] = draw_box(draw, section_font, body_font, box["title"], box["lines"], *box["xy"])

    arrow(draw, (400, 350), (400, 510))
    arrow(draw, (840, 350), (520, 510))
    arrow(draw, (680, 350), (860, 510))
    arrow(draw, (400, 350), (860, 510))

    draw.text((250, 430), "1 to many", fill="#7a5412", font=section_font)
    draw.text((615, 430), "1 to many", fill="#7a5412", font=section_font)
    draw.text((920, 430), "1 to many", fill="#7a5412", font=section_font)
    draw.text((730, 455), "1 to many", fill="#7a5412", font=section_font)

    legend_x = 1080
    legend_y = 640
    draw.rounded_rectangle((legend_x, legend_y, legend_x + 280, legend_y + 150), radius=18, fill="white", outline="#1c2b2d", width=2)
    draw.text((legend_x + 18, legend_y + 16), "Legend", fill="#143128", font=section_font)
    draw.text((legend_x + 18, legend_y + 54), "PK = Primary Key", fill="#1c2b2d", font=body_font)
    draw.text((legend_x + 18, legend_y + 86), "FK = Foreign Key", fill="#1c2b2d", font=body_font)
    draw.text((legend_x + 18, legend_y + 118), "Announcements is standalone", fill="#1c2b2d", font=body_font)

    image.save(JPG_PATH, quality=95)


def generate_erd_pdf():
    c = canvas.Canvas(str(PDF_PATH), pagesize=landscape(TABLOID))
    width, height = landscape(TABLOID)
    c.setFillColor(colors.HexColor("#f6f1e8"))
    c.rect(0, 0, width, height, fill=1, stroke=0)
    c.setFillColor(colors.HexColor("#143128"))
    c.setFont("Helvetica-Bold", 24)
    c.drawString(40, height - 40, "BarangayConnect Final Entity Relationship Diagram")
    c.setFillColor(colors.HexColor("#5e6e68"))
    c.setFont("Helvetica", 13)
    c.drawString(42, height - 62, "Local government service portal using an embedded SQLite database")

    pdf_boxes = [
        ("Residents", ["PK ResidentId", "FullName", "HouseholdNo", "ContactNumber", "EmailAddress", "Purok", "RegisteredOn"], 35, 255),
        ("Services", ["PK ServiceId", "Name", "Office", "Description", "Schedule", "Requirements"], 330, 255),
        ("Announcements", ["PK AnnouncementId", "Title", "Category", "Summary", "PublishedOn", "Audience"], 625, 255),
        ("Appointments", ["PK AppointmentId", "FK ResidentId", "FK ServiceId", "AppointmentDate", "TimeSlot", "Status", "Notes"], 120, 40),
        ("ServiceRequests", ["PK RequestId", "FK ResidentId", "FK ServiceId", "Description", "Priority", "Status", "SubmittedOn"], 450, 40),
    ]

    def pdf_box(title, lines, x, y):
        c.setStrokeColor(colors.HexColor("#1c2b2d"))
        c.setFillColor(colors.white)
        c.roundRect(x, y, 220, 175, 12, fill=1, stroke=1)
        c.setFillColor(colors.HexColor("#0c7c59"))
        c.roundRect(x, y + 141, 220, 34, 12, fill=1, stroke=0)
        c.setFillColor(colors.white)
        c.setFont("Helvetica-Bold", 13)
        c.drawString(x + 10, y + 153, title)
        c.setFillColor(colors.HexColor("#1c2b2d"))
        c.setFont("Helvetica", 11)
        current_y = y + 126
        for line in lines:
            c.drawString(x + 10, current_y, line)
            current_y -= 16

    for title, lines, x, y in pdf_boxes:
        pdf_box(title, lines, x, y)

    c.setStrokeColor(colors.HexColor("#7a5412"))
    c.setLineWidth(2)
    c.line(145, 255, 145, 215)
    c.line(445, 255, 320, 215)
    c.line(440, 255, 560, 215)
    c.line(740, 255, 560, 215)
    c.setFont("Helvetica", 11)
    c.setFillColor(colors.HexColor("#7a5412"))
    c.drawString(115, 220, "1:M")
    c.drawString(355, 220, "1:M")
    c.drawString(515, 220, "1:M")
    c.drawString(650, 220, "1:M")

    c.showPage()
    c.save()


def add_title(slide, title, subtitle=""):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(8.5), Inches(1.0))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_run = title_frame.paragraphs[0].runs[0]
    title_run.font.size = Pt(26)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(20, 49, 40)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.1), Inches(8.8), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_run = subtitle_frame.paragraphs[0].runs[0]
        subtitle_run.font.size = Pt(12)
        subtitle_run.font.color.rgb = RGBColor(94, 110, 104)


def add_bullets(slide, items, left=0.8, top=1.8, width=8.5, height=4.8):
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = textbox.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(28, 43, 45)


def generate_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "BarangayConnect", "IT114L Final Project Presentation")
    add_bullets(
        slide,
        [
            "Institutional website for a barangay or local government office",
            "Built with ASP.NET Core MVC on .NET 10",
            "Uses an embedded SQLite database stored in the project source",
        ],
        top=1.9,
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Problem and Objective", "Why the system was created")
    add_bullets(
        slide,
        [
            "Manual service tracking can slow down resident transactions and office coordination.",
            "The project centralizes announcements, resident records, appointments, and requests.",
            "The goal is to demonstrate web development, UI design, and SQL-backed data management.",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "System Modules", "Main features of the web application")
    add_bullets(
        slide,
        [
            "Dashboard with service statistics and recent activity",
            "Announcements page for public notices",
            "Resident registry with data entry form",
            "Service catalog with office schedules and requirements",
            "Appointment scheduling form and listing",
            "Service request submission and tracking",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Database Design", "Embedded SQLite for easy updates")
    add_bullets(
        slide,
        [
            "Tables: Residents, Services, Announcements, Appointments, ServiceRequests",
            "Appointments and ServiceRequests use foreign keys to connect Residents and Services",
            "Database file is stored locally in App_Data/barangayconnect.db",
            "Raw SQL statements are used for transparent retrieval and insert operations",
        ],
    )
    slide.shapes.add_picture(str(JPG_PATH), Inches(8.7), Inches(1.2), width=Inches(4.0))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Technology Stack", "Frameworks and implementation choices")
    add_bullets(
        slide,
        [
            ".NET 10 and ASP.NET Core MVC",
            "Microsoft.Data.Sqlite for in-house SQL database handling",
            "Bootstrap-based responsive user interface",
            "Repository pattern for database initialization and data access",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Conclusion", "Project outcome")
    add_bullets(
        slide,
        [
            "BarangayConnect satisfies the requirement for multiple pages, forms, navigation, and database integration.",
            "The SQLite approach keeps the project easy to deploy, inspect, and modify.",
            "The application demonstrates a practical institutional system that can be extended further with authentication and reporting.",
        ],
    )

    prs.save(PPTX_PATH)


if __name__ == "__main__":
    DOCS.mkdir(parents=True, exist_ok=True)
    generate_erd_jpg()
    generate_erd_pdf()
    generate_presentation()
